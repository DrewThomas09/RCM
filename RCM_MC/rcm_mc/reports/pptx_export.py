"""
Step 93: PowerPoint export.

Generates a 5-slide IC deck from simulation results.
Requires python-pptx: pip install python-pptx
"""
from __future__ import annotations

import os
from typing import Any, Dict, Optional

import pandas as pd

from ..infra.logger import logger


def _fmt_money(val: float) -> str:
    if abs(val) >= 1_000_000:
        return f"${val / 1_000_000:,.1f}M"
    return f"${val / 1_000:,.0f}K"


def generate_pptx(
    outdir: str,
    *,
    hospital_name: str = "Target Hospital",
    ev_multiple: float = 8.0,
    annual_revenue: float = 0.0,
    n_sims: int = 30000,
) -> Optional[str]:
    """Generate a PowerPoint deck from simulation outputs."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        logger.warning("python-pptx not installed; skipping PPTX export. Run: pip install python-pptx")
        return None

    summary_path = os.path.join(outdir, "summary.csv")
    if not os.path.exists(summary_path):
        logger.warning("summary.csv not found; cannot generate PPTX")
        return None

    summary = pd.read_csv(summary_path, index_col=0)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"RCM Opportunity Analysis: {hospital_name}"
    slide.placeholders[1].text = (
        f"Monte Carlo Simulation ({n_sims:,} iterations)\n"
        f"Annual Revenue: {_fmt_money(annual_revenue)}"
    )

    # Slide 2: Headline numbers
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Opportunity Summary"
    if "ebitda_drag" in summary.index:
        ebitda = float(summary.loc["ebitda_drag", "mean"])
        p10 = float(summary.loc["ebitda_drag", "p10"])
        p90 = float(summary.loc["ebitda_drag", "p90"])
        ev = ebitda * ev_multiple
        body = slide.placeholders[1]
        body.text = (
            f"Annual EBITDA Opportunity: {_fmt_money(ebitda)}\n"
            f"Range: {_fmt_money(p10)} (P10) to {_fmt_money(p90)} (P90)\n"
            f"Enterprise Value @ {ev_multiple:.0f}x: {_fmt_money(ev)}"
        )

    # Slide 3: Value Drivers (from sensitivity.csv if available)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top Value Drivers"
    sens_path = os.path.join(outdir, "sensitivity.csv")
    driver_lines = []
    if os.path.exists(sens_path):
        df_sens = pd.read_csv(sens_path)
        label_col = "driver_label" if "driver_label" in df_sens.columns else "driver"
        for _, row in df_sens.head(5).iterrows():
            lbl = str(row.get(label_col, row.get("driver", "")))
            corr = float(row.get("corr", 0))
            driver_lines.append(f"{lbl}: {corr:.2f} correlation to EBITDA drag")
    if driver_lines:
        slide.placeholders[1].text = "\n".join(driver_lines)
    else:
        slide.placeholders[1].text = "Sensitivity data not available. Run with --stress for driver analysis."

    # Slide 4: Risk Register (data-driven from simulation outputs)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Risks"
    risk_lines = []
    if "ebitda_drag" in summary.index:
        p10 = float(summary.loc["ebitda_drag", "p10"])
        p90 = float(summary.loc["ebitda_drag", "p90"])
        spread = p90 - p10
        risk_lines.append(f"Tail Risk: P90 exceeds P10 by {_fmt_money(spread)}, indicating significant outcome uncertainty")
    for m, label in [("drag_denial_writeoff", "Denial Write-Off"), ("drag_dar_total", "A/R Days")]:
        if m in summary.index:
            val_p90 = float(summary.loc[m, "p90"])
            val_mean = float(summary.loc[m, "mean"])
            if val_p90 > val_mean * 1.4:
                risk_lines.append(f"{label} Volatility: P90 is {val_p90/val_mean:.0%} of mean, suggesting high variability")
    risk_lines.append("Data Quality: Results depend on accuracy of denial rate and write-off assumptions")
    risk_lines.append("Payer Policy: Commercial payer tightening could increase denial rates beyond modeled levels")
    slide.placeholders[1].text = "\n".join(risk_lines[:5])

    # Slide 5: Recommended Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommended Next Steps"
    next_steps = [
        "1. Request 12-24 months of claims-level denial data from management",
        "2. Validate IDR and FWR assumptions against actual remittance data",
        "3. Assess denial management staffing capacity and backlog",
        "4. Model CDI and prior-auth automation ROI for 100-day plan",
        "5. Refine underwriting credit based on calibrated results",
    ]
    slide.placeholders[1].text = "\n".join(next_steps)

    pptx_path = os.path.join(outdir, "report.pptx")
    prs.save(pptx_path)
    logger.info("Wrote: %s", pptx_path)
    return pptx_path
